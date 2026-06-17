from __future__ import annotations

import csv
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[1]
RESULTS_ROOT = PROJECT_ROOT / "Results"
REPORTS_DIR = PROJECT_ROOT / "Reports"
REPORTS_DIR.mkdir(exist_ok=True)

TITLE = "政策性調薪是否縮減了職場貧富差距？"
SUBTITLE = "2016-2026 年主管職與基層勞工薪資倍數變動分析"
AUTHOR = "414351346 侯卉倢"
COURSE = "資料科學與經濟分析入門：期末專案 Part 3"
PROFESSOR = "指導教授：朱建達 (Jian-Da Zhu) 老師"


def latest_results_dir() -> Path:
    dated = [p for p in RESULTS_ROOT.iterdir() if p.is_dir() and len(p.name) == 10]
    if not dated:
        raise FileNotFoundError("No dated Results/YYYY-MM-DD directory found.")
    return sorted(dated, key=lambda p: p.name)[-1]


def read_csv_rows(path: Path) -> list[dict[str, str]]:
    with path.open("r", encoding="utf-8-sig", newline="") as f:
        return list(csv.DictReader(f))


def fmt_money(value: str | float) -> str:
    return f"{float(value):,.0f}"


def fmt_num(value: str | float, digits: int = 3) -> str:
    return f"{float(value):,.{digits}f}"


def fmt_pct_or_dash(value: str) -> str:
    if value in ("", "NA", "NaN", None):
        return "-"
    return fmt_num(value, 2) + "%"


def value_from_summary(rows: list[dict[str, str]], item: str) -> str:
    for row in rows:
        if row["item"] == item:
            return row["value"]
    raise KeyError(item)


def set_run_font(paragraph, size=20, bold=False, color=RGBColor(30, 42, 56)):
    for run in paragraph.runs:
        run.font.name = "Microsoft JhengHei"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def add_title(slide, title: str, subtitle: str | None = None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(12.25), Inches(0.65))
    p = box.text_frame.paragraphs[0]
    p.text = title
    set_run_font(p, size=24, bold=True, color=RGBColor(24, 50, 80))
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(0.9), Inches(12.0), Inches(0.35))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        set_run_font(sp, size=12, color=RGBColor(100, 100, 100))


def add_bullets(slide, bullets: list[str], x=0.8, y=1.35, w=11.8, h=5.3, size=18):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, text in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.space_after = Pt(8)
        set_run_font(p, size=size)


def add_footer(slide, results_date: str):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(12.0), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.text = f"{COURSE}｜{AUTHOR}｜資料版本：{results_date}"
    p.alignment = PP_ALIGN.RIGHT
    set_run_font(p, size=8, color=RGBColor(120, 120, 120))


def add_table(slide, rows: list[list[str]], x, y, w, h, font_size=10):
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h)).table
    for r_idx, row in enumerate(rows):
        for c_idx, text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT
                set_run_font(
                    paragraph,
                    size=font_size,
                    bold=(r_idx == 0),
                    color=RGBColor(255, 255, 255) if r_idx == 0 else RGBColor(30, 42, 56),
                )
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(34, 91, 122) if r_idx == 0 else RGBColor(245, 248, 250)
    return table


def add_image(slide, path: Path, x, y, w, h):
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def main() -> Path:
    results_dir = latest_results_dir()
    results_date = results_dir.name

    summary = read_csv_rows(results_dir / "policy_wage_gap_summary.csv")
    stats = read_csv_rows(results_dir / "policy_wage_gap_descriptive_stats.csv")
    yearly = read_csv_rows(results_dir / "policy_wage_gap_yearly.csv")

    requested_period = value_from_summary(summary, "requested_period")
    available_period = value_from_summary(summary, "available_empirical_period")
    ratio_first = value_from_summary(summary, "ratio_first_year")
    ratio_last = value_from_summary(summary, "ratio_last_year")
    ratio_change = value_from_summary(summary, "ratio_change")
    gap_first = value_from_summary(summary, "gap_first_year")
    gap_last = value_from_summary(summary, "gap_last_year")
    gap_change = value_from_summary(summary, "gap_change")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(237, 244, 248)
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(12.0), Inches(1.0))
    p = title_box.text_frame.paragraphs[0]
    p.text = TITLE
    p.alignment = PP_ALIGN.CENTER
    set_run_font(p, size=32, bold=True, color=RGBColor(24, 50, 80))
    sub_box = slide.shapes.add_textbox(Inches(1.1), Inches(2.35), Inches(11.2), Inches(0.8))
    sp = sub_box.text_frame.paragraphs[0]
    sp.text = SUBTITLE
    sp.alignment = PP_ALIGN.CENTER
    set_run_font(sp, size=22, bold=True, color=RGBColor(40, 80, 120))
    add_bullets(
        slide,
        [COURSE, PROFESSOR, AUTHOR, f"資料版本：{results_date}"],
        x=2.25,
        y=4.1,
        w=8.9,
        h=1.6,
        size=16,
    )

    # 2. Research question
    slide = prs.slides.add_slide(blank)
    add_title(slide, "研究問題與核心結論")
    add_bullets(
        slide,
        [
            f"研究問題：政策性調薪是否讓主管職與基層勞工之間的薪資倍數下降？",
            f"原題設定期間為 {requested_period}；目前官方薪資原始檔可支持 {available_period} 的可重現分析。",
            f"主管職/基層勞工薪資倍數由 {ratio_first} 倍降至 {ratio_last} 倍，變化 {ratio_change} 倍。",
            f"月薪差距由 {fmt_money(gap_first)} 元變為 {fmt_money(gap_last)} 元，差距變化 {fmt_money(gap_change)} 元。",
            "初步判讀：倍數下降，表示相對薪資差距縮小；但絕對金額差距幾乎持平。",
        ],
    )
    add_footer(slide, results_date)

    # 3. Data and definitions
    slide = prs.slides.add_slide(blank)
    add_title(slide, "資料來源與職類定義")
    add_bullets(
        slide,
        [
            "資料來源：行政院主計總處「工業及服務業廠商僱用按月、按日及按時計薪者平均最低薪資」。",
            "薪資指標：按月計薪者每人每月平均最低薪資，單位為新臺幣元。",
            "主管職：原始職類「主管及監督人員」。",
            "基層勞工：原始職類「基層技術工及勞力工」。",
            "注意：本題使用名目月薪計算薪資倍數，不以 CPI 或通膨調整後實質薪資為核心。",
        ],
    )
    add_footer(slide, results_date)

    # 4. Cleaning workflow
    slide = prs.slides.add_slide(blank)
    add_title(slide, "資料清理與可重現流程")
    add_bullets(
        slide,
        [
            "RawData：保存未修改原始 CSV/XML；所有清理由 R 腳本自動執行。",
            "清理前檢視：程式輸出 str/head/tail/summary/colnames/nrow/ncol 至 raw_data_inspection_report.txt。",
            "缺失值處理：原始 '-' 與空值轉為 NA；月薪缺失或非正值列另存 excluded_wage_records。", 
            "年份核對：以主計總處民國 104-113 年表 5 的總計月薪驗證 2015-2024 對照。",
            "版本管理：結果輸出至 Results/YYYY-MM-DD/，本次使用 2026-06-17。",
        ],
    )
    add_footer(slide, results_date)

    # 5. Descriptive statistics table
    slide = prs.slides.add_slide(blank)
    add_title(slide, "描述統計：N、Mean、Median、Min、Max、SD")
    table_rows = [["變數", "N", "Mean", "Median", "Min", "Max", "SD"]]
    for row in stats:
        table_rows.append([
            row["variable"],
            row["N"],
            fmt_num(row["Mean"], 3),
            fmt_num(row["Median"], 3),
            fmt_num(row["Min"], 3),
            fmt_num(row["Max"], 3),
            fmt_num(row["SD"], 3),
        ])
    add_table(slide, table_rows, 0.35, 1.25, 12.65, 4.45, font_size=8.5)
    add_footer(slide, results_date)

    # 6. Descriptive statistics chart
    slide = prs.slides.add_slide(blank)
    add_title(slide, "描述統計圖：平均值比較")
    add_image(slide, results_dir / "03_policy_wage_gap_descriptive_stats.png", 0.85, 1.05, 11.7, 5.55)
    add_footer(slide, results_date)

    # 7. Wage trend chart
    slide = prs.slides.add_slide(blank)
    add_title(slide, "主管職與基層勞工名目月薪走勢")
    add_image(slide, results_dir / "01_supervisor_grassroots_wage_trend.png", 0.85, 1.05, 11.7, 5.55)
    add_footer(slide, results_date)

    # 8. Wage ratio chart
    slide = prs.slides.add_slide(blank)
    add_title(slide, "主管職／基層勞工薪資倍數變動")
    add_image(slide, results_dir / "02_supervisor_grassroots_wage_ratio.png", 0.85, 1.05, 11.7, 5.55)
    add_footer(slide, results_date)

    # 9. Yearly data table
    slide = prs.slides.add_slide(blank)
    add_title(slide, "年度薪資倍數資料表")
    yearly_rows = [["年", "主管職月薪", "基層月薪", "差距", "倍數", "主管年增率", "基層年增率"]]
    for row in yearly:
        yearly_rows.append([
            row["year"],
            fmt_money(row["supervisor_wage"]),
            fmt_money(row["grassroots_wage"]),
            fmt_money(row["wage_gap"]),
            fmt_num(row["wage_ratio"], 3),
            fmt_pct_or_dash(row["supervisor_growth"]),
            fmt_pct_or_dash(row["grassroots_growth"]),
        ])
    add_table(slide, yearly_rows, 0.55, 1.2, 12.2, 4.95, font_size=8.5)
    add_footer(slide, results_date)

    # 10. Interpretation
    slide = prs.slides.add_slide(blank)
    add_title(slide, "分析解讀")
    add_bullets(
        slide,
        [
            f"相對差距：薪資倍數由 {ratio_first} 降至 {ratio_last}，代表主管職相對於基層勞工的薪資優勢降低。",
            f"絕對差距：月薪差距只從 {fmt_money(gap_first)} 元變為 {fmt_money(gap_last)} 元，幾乎沒有明顯縮小。",
            "基層勞工平均年增率高於主管職，推動倍數下降。",
            "若政策性調薪主要作用於低薪族群，從倍數角度看有縮小相對差距的跡象。",
            "但因資料為平均最低薪資，結論應界定為「僱用門檻薪資差距」，不可直接推論所有在職者實領薪資差距。",
        ],
    )
    add_footer(slide, results_date)

    # 11. Limitations
    slide = prs.slides.add_slide(blank)
    add_title(slide, "限制與後續延伸")
    add_bullets(
        slide,
        [
            "期間限制：目前原始薪資資料支援 2016-2024；2025-2026 需補官方表格後更新。",
            "指標限制：資料為廠商僱用平均最低薪資，不是全體受僱者實領平均薪資。",
            "分類限制：主管職與基層勞工以職類對應，未控制產業、地區、企業規模與工時差異。",
            "政策歸因限制：倍數變化與政策性調薪方向一致，但仍需結合最低工資政策時點與其他總體因素才能做因果判斷。",
            "延伸方向：補齊 2025-2026、加入企業規模分組、比較其他低薪職類或加入政策事件線。",
        ],
    )
    add_footer(slide, results_date)

    output_path = REPORTS_DIR / f"policy_wage_gap_report_{results_date}.pptx"
    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    path = main()
    print(path)
